"""Build the empty shell deck for the Engineering Automation Opportunities meeting.
Title slide, agenda slide, and a 3-slide shell for each of the 3 topics — titles only, no content yet."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Duetto brand palette ──────────────────────────────────────────────────────
MIDNIGHT = RGBColor(0x0E, 0x21, 0x24)
LUCENT   = RGBColor(0xC4, 0xFF, 0x45)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
INK      = RGBColor(0x0E, 0x21, 0x24)
MUTED    = RGBColor(0x5B, 0x6B, 0x6B)
CARD_BG  = RGBColor(0xF3, 0xF6, 0xF2)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

FONT_HEAD = "Arial"
FONT_BODY = "Arial"

# Topic order per this meeting's agenda — topic 3 is DQE
TOPICS = ["Auto-Provisioning", "Migrations", "Data Quality Engine (DQE)"]


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def no_line(shape):
    shape.line.fill.background()


def txt(slide, x, y, w, h, text, size=14, color=INK, bold=False, italic=False,
        font=FONT_BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    return box


def bullets(slide, x, y, w, h, items, size=14, color=INK, font=FONT_BODY,
            space_after=8, line_spacing=1.15, bullet_color=None, numbered=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    bc = bullet_color or color
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        pPr = p._p.get_or_add_pPr()
        if numbered:
            buFont = pPr.makeelement(qn('a:buFont'), {'typeface': font})
            buAutoNum = pPr.makeelement(qn('a:buAutoNum'), {'type': 'arabicPeriod'})
            buClr = pPr.makeelement(qn('a:buClr'), {})
            srgb = pPr.makeelement(qn('a:srgbClr'), {'val': '%02X%02X%02X' % (bc[0], bc[1], bc[2])})
            buClr.append(srgb)
            pPr.append(buClr)
            pPr.append(buFont)
            pPr.append(buAutoNum)
        else:
            buChar = pPr.makeelement(qn('a:buChar'), {'char': '–'})
            buFont = pPr.makeelement(qn('a:buFont'), {'typeface': font})
            buClr = pPr.makeelement(qn('a:buClr'), {})
            srgb = pPr.makeelement(qn('a:srgbClr'), {'val': '%02X%02X%02X' % (bc[0], bc[1], bc[2])})
            buClr.append(srgb)
            pPr.append(buClr)
            pPr.append(buFont)
            pPr.append(buChar)
        pPr.set('indent', '-228600')
        pPr.set('marL', '228600')
        r = p.add_run()
        r.text = item
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = font
    return box


def header(slide, title, kicker):
    set_bg(slide, WHITE)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1.55))
    band.fill.solid(); band.fill.fore_color.rgb = MIDNIGHT
    no_line(band)
    txt(slide, 0.6, 0.4, 11.5, 0.3, kicker, size=10.5, color=LUCENT, bold=True, font=FONT_BODY)
    txt(slide, 0.6, 0.72, 12.2, 0.6, title, size=30, color=WHITE, bold=True, font=FONT_HEAD)


# ══════════════════════════════════════════════════════════════════════════
# TITLE SLIDE
# ══════════════════════════════════════════════════════════════════════════
s0 = add_slide()
set_bg(s0, MIDNIGHT)

big_circle = s0.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.6), Inches(-2.4), Inches(6.5), Inches(6.5))
big_circle.fill.solid(); big_circle.fill.fore_color.rgb = LUCENT
no_line(big_circle)
srgb_el = big_circle.fill.fore_color._xFill.find(qn('a:srgbClr'))
alpha_el = srgb_el.makeelement(qn('a:alpha'), {'val': '12000'})
srgb_el.append(alpha_el)

txt(s0, 0.9, 2.4, 11.0, 1.4, "Engineering Automation\nOpportunities", size=44, color=WHITE,
    bold=True, font=FONT_HEAD, line_spacing=1.1)
txt(s0, 0.9, 4.15, 10.5, 0.5, "Prep for Engineering", size=16, color=LUCENT, italic=True, font=FONT_BODY)


# ══════════════════════════════════════════════════════════════════════════
# AGENDA SLIDE
# ══════════════════════════════════════════════════════════════════════════
s1 = add_slide()
header(s1, "Agenda", "ENGINEERING AUTOMATION OPPORTUNITIES")

card_y, card_h, gap, card_w = 2.35, 1.35, 0.35, 12.1
for i, topic in enumerate(TOPICS):
    y = card_y + i * (card_h + gap)
    c = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(y), Inches(card_w), Inches(card_h))
    c.adjustments[0] = 0.08
    c.fill.solid(); c.fill.fore_color.rgb = CARD_BG
    no_line(c)

    num = s1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95), Inches(y + card_h/2 - 0.32), Inches(0.64), Inches(0.64))
    num.fill.solid(); num.fill.fore_color.rgb = LUCENT
    no_line(num)
    tf = num.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(i + 1)
    r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = MIDNIGHT; r.font.name = FONT_BODY

    txt(s1, 2.0, y, card_w - 1.6, card_h, topic, size=20, color=MIDNIGHT, bold=True,
        font=FONT_BODY, anchor=MSO_ANCHOR.MIDDLE)


# ══════════════════════════════════════════════════════════════════════════
# 3-SLIDE SHELL PER TOPIC — titles only, no content yet
# ══════════════════════════════════════════════════════════════════════════
for topic_idx, topic in enumerate(TOPICS, start=1):
    for slide_num in range(1, 4):
        s = add_slide()
        header(s, topic, f"AGENDA ITEM {topic_idx} OF {len(TOPICS)}  ·  SLIDE {slide_num} OF 3")

prs.save("/Users/tylerlinton/DQE/Engineering_Automation_Opportunities_Shell.pptx")
print("saved")
