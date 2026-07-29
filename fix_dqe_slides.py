"""Reformat slides 4-6 (DQE section) to match the visual format of slides 8-10
(Auto-Provisioning / Migration): plain white background, single Sora title line
in the 'Topic | Subtitle' pattern, small Duetto logo icon bottom-left — instead
of the navy header band + two-line Arial title used before."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

SRC  = "/Users/tylerlinton/Downloads/_Onboarding automation opportunities.pptx"
DST  = "/Users/tylerlinton/Downloads/_Onboarding automation opportunities (updated).pptx"
LOGO = "/tmp/duetto_logo.png"

BLACK = RGBColor(0x00, 0x00, 0x00)

prs = Presentation(SRC)

new_titles = {
    3: "DQE | Current Workflow",
    4: "DQE | Where the Tool Fits In",
    5: "DQE | The Engineering Ask",
}

# Position/size of the title textbox and logo on the reference slides (8, 9, 10)
TITLE_X, TITLE_Y, TITLE_W, TITLE_H = 228600, 228600, 8686800, 669600
LOGO_X, LOGO_Y, LOGO_W, LOGO_H = 198823, 4800600, 148850, 117025

for idx, title_text in new_titles.items():
    slide = prs.slides[idx]
    shapes = list(slide.shapes)

    # Remove the first 4 shapes: navy band, kicker line, old title, old subtitle
    for sh in shapes[:4]:
        sh._element.getparent().remove(sh._element)

    # New single-line Sora title, matching slides 8-10 exactly
    box = slide.shapes.add_textbox(Emu(TITLE_X), Emu(TITLE_Y), Emu(TITLE_W), Emu(TITLE_H))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = 0.8
    r = p.add_run()
    r.text = title_text
    r.font.size = Pt(24)
    r.font.name = "Sora"
    r.font.color.rgb = BLACK

    # Small Duetto logo icon, bottom-left — same spot as slides 7-9
    slide.shapes.add_picture(LOGO, Emu(LOGO_X), Emu(LOGO_Y), Emu(LOGO_W), Emu(LOGO_H))

prs.save(DST)
print("saved:", DST)
