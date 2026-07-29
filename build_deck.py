"""Build the DQE section (3 slides) for the Onboarding Automation Opportunities deck.
Styled to Duetto brand guidelines: Midnight Green / Lucent Green palette, Sora (fallback Arial),
70/30 Midnight-to-Lucent weighting, no accent stripes."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Duetto brand palette ──────────────────────────────────────────────────────
MIDNIGHT   = RGBColor(0x0E, 0x21, 0x24)   # Midnight Green — primary, ~70% weight
MIDNIGHT_2 = RGBColor(0x16, 0x2E, 0x32)   # slightly lighter midnight for tag chips
LUCENT     = RGBColor(0xC4, 0xFF, 0x45)   # Lucent Green — primary accent, ~30% weight
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
INK        = RGBColor(0x0E, 0x21, 0x24)   # body text = Midnight Green
MUTED      = RGBColor(0x5B, 0x6B, 0x6B)
CARD_BG    = RGBColor(0xF3, 0xF6, 0xF2)   # faint green-grey neutral
AMBER      = RGBColor(0xB2, 0x77, 0x0E)   # semantic "still manual" flag — used sparingly
AMBER_BG   = RGBColor(0xF8, 0xEF, 0xDC)
GREEN_TXT  = RGBColor(0x2F, 0x5A, 0x1E)   # readable dark-green text for "automated" pill

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# Sora is Duetto's primary typeface; Arial is the brand-approved fallback for
# environments where Sora isn't installed (safer for a deck presented off-machine).
FONT_HEAD = "Arial"
FONT_BODY = "Arial"


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def no_line(shape):
    shape.line.fill.background()


def txt(slide, x, y, w, h, text, size=14, color=INK, bold=False, italic=False,
        font=FONT_BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    return box


def bullets(slide, x, y, w, h, items, size=14, color=INK, font=FONT_BODY,
            space_after=8, line_spacing=1.12, bullet_color=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    bc = bullet_color or color
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        pPr = p._p.get_or_add_pPr()
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '–'})
        buFont = pPr.makeelement(qn('a:buFont'), {'typeface': font})
        buClr = pPr.makeelement(qn('a:buClr'), {})
        srgb = pPr.makeelement(qn('a:srgbClr'), {'val': '%02X%02X%02X' % (bc[0], bc[1], bc[2])})
        buClr.append(srgb)
        pPr.append(buClr)
        pPr.append(buFont)
        pPr.append(buChar)
        pPr.set('indent', '-182880')
        pPr.set('marL', '182880')
        r = p.add_run()
        r.text = item
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = font
    return box


def card(slide, x, y, w, h, fill=CARD_BG, shadow=True, radius=0.06):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.adjustments[0] = radius
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    no_line(shp)
    if shadow:
        el = shp._element.spPr
        effectLst = el.makeelement(qn('a:effectLst'), {})
        outerShdw = el.makeelement(qn('a:outerShdw'), {
            'blurRad': '90000', 'dist': '25000', 'dir': '5400000', 'rotWithShape': '0'
        })
        clr = el.makeelement(qn('a:srgbClr'), {'val': '0E2124'})
        alpha = el.makeelement(qn('a:alpha'), {'val': '22000'})
        clr.append(alpha)
        outerShdw.append(clr)
        effectLst.append(outerShdw)
        el.append(effectLst)
    return shp


def header(slide, title, subtitle, kicker="DQE  ·  ONBOARDING AUTOMATION OPPORTUNITIES"):
    set_bg(slide, WHITE)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1.55))
    band.fill.solid(); band.fill.fore_color.rgb = MIDNIGHT
    no_line(band)
    txt(slide, 0.6, 0.22, 11.5, 0.3, kicker, size=10.5, color=LUCENT, bold=True, font=FONT_BODY)
    txt(slide, 0.6, 0.52, 12.2, 0.6, title, size=28, color=WHITE, bold=True, font=FONT_HEAD)
    txt(slide, 0.6, 1.63, 12.2, 0.42, subtitle, size=14, color=MUTED, italic=True, font=FONT_BODY)


def circle_num(slide, x, y, d, num):
    """Lucent Green accent circle — the signature Duetto accent moment on a light card."""
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    c.fill.solid(); c.fill.fore_color.rgb = LUCENT
    no_line(c)
    tf = c.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(num)
    r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = MIDNIGHT; r.font.name = FONT_BODY
    return c


# ══════════════════════════════════════════════════════════════════════════
# TITLE SLIDE
# ══════════════════════════════════════════════════════════════════════════
s0 = add_slide()
set_bg(s0, MIDNIGHT)

# Large soft Lucent circle, off-canvas top-right — a quiet brand accent moment,
# echoing the circular motif used for numbering throughout the deck
big_circle = s0.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.6), Inches(-2.4), Inches(6.5), Inches(6.5))
big_circle.fill.solid(); big_circle.fill.fore_color.rgb = LUCENT
no_line(big_circle)
srgb_el = big_circle.fill.fore_color._xFill.find(qn('a:srgbClr'))
alpha_el = srgb_el.makeelement(qn('a:alpha'), {'val': '12000'})
srgb_el.append(alpha_el)

# DQE badge — small rounded chip echoing the app's nav badge
badge = s0.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(0.85), Inches(1.15), Inches(0.5))
badge.adjustments[0] = 0.4
badge.fill.solid(); badge.fill.fore_color.rgb = LUCENT
no_line(badge)
tf = badge.text_frame
tf.margin_left = tf.margin_right = 0
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "DQE"
r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = MIDNIGHT; r.font.name = FONT_BODY

# Eyebrow / kicker
txt(s0, 0.9, 1.55, 8.0, 0.35, "ONBOARDING AUTOMATION OPPORTUNITIES  ·  SECTION 1 OF 3", size=13, color=LUCENT,
    bold=True, font=FONT_BODY)

# Title
txt(s0, 0.85, 2.05, 11.6, 1.7, "Data Quality Engine (DQE)", size=46, color=WHITE,
    bold=True, font=FONT_HEAD, line_spacing=1.05)

# Supporting line
txt(s0, 0.9, 3.35, 10.8, 0.9,
    "Where the tool stands today, what it already automates, and the API access "
    "that DQE needs — alongside Auto-Provisioning and Migrations — to go further.",
    size=16, color=WHITE, italic=True, font=FONT_BODY, line_spacing=1.25)

# Presenter / date footer block
txt(s0, 0.9, 6.55, 9.0, 0.4, "Tyler Linton  ·  Prep for Engineering — Automation Opportunities",
    size=12.5, color=WHITE, font=FONT_BODY)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — DQE Today: The Current Workflow
# ══════════════════════════════════════════════════════════════════════════
s1 = add_slide()
header(s1, "DQE Today: The Current Workflow",
       "Every discrepancy investigation today is a manual, report-driven process.")

steps = [
    "Customer reports a data discrepancy by opening a support case",
    "Analyst reviews the DVA to determine whether Duetto is overstating or understating revenue/occupancy",
    "Based on the variance, the analyst follows the appropriate investigation path",
    "Analyst manually compares reservations between Duetto and the PMS",
    "Additional PMS reports are exported as needed to identify specific issues",
    "Root cause is identified and the appropriate corrective action is taken",
]

flow_x, flow_w = 0.6, 7.55
row_y, row_h, gap = 2.28, 0.735, 0.09
for i, step in enumerate(steps):
    y = row_y + i * (row_h + gap)
    card(s1, flow_x, y, flow_w, row_h, fill=CARD_BG, shadow=False, radius=0.12)
    circle_num(s1, flow_x + 0.18, y + row_h/2 - 0.22, 0.44, i + 1)
    txt(s1, flow_x + 0.8, y, flow_w - 1.0, row_h, step, size=12.5, color=INK,
        font=FONT_BODY, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)

# Side callout — investigation paths (Midnight card, Lucent accents — signature contrast panel)
side_x, side_w = 8.45, 4.28
card(s1, side_x, 2.28, side_w, 4.68, fill=MIDNIGHT, shadow=True, radius=0.05)
txt(s1, side_x + 0.35, 2.55, side_w - 0.7, 0.3, "INVESTIGATION PATHS", size=12, color=LUCENT, bold=True, font=FONT_BODY)
txt(s1, side_x + 0.35, 2.88, side_w - 0.7, 0.5, "What additional PMS exports typically uncover:",
    size=11.5, color=WHITE, italic=True, font=FONT_BODY, line_spacing=1.1)
paths = [
    "Missing reservations",
    "Rate discrepancies",
    "Source code issues",
    "Package, tax, or gross vs. net differences",
    "Other configuration or data issues",
]
bullets(s1, side_x + 0.35, 3.55, side_w - 0.7, 3.2, paths, size=13, color=WHITE,
        font=FONT_BODY, space_after=14, line_spacing=1.15, bullet_color=LUCENT)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Where the DQE Tool Fits In
# ══════════════════════════════════════════════════════════════════════════
s2 = add_slide()
header(s2, "Where the DQE Tool Fits In",
       "The tool automates the analysis step — but it still depends entirely on manually exported files to run.")

items2 = [
    ("Automated root-cause classification. ",
     "Identifies rooms/revenue variance and classifies WHY it happened — not just that it happened — using a fixed diagnostic rule set applied consistently across every analyst."),
    ("Automated reservation-level comparison. ",
     "Cross-references Duetto against the PMS automatically, including a new Sync Gap detection capability that determines whether a discrepancy is a true Duetto data issue or an upstream PMS sync/integration issue."),
    ("Automated, structured output. ",
     "Generates a report, an Excel export, and a Monday.com audit trail automatically — replacing ad hoc analyst notes with a consistent, reviewable record."),
]

card_y, card_h, card_gap, card_w = 2.35, 1.28, 0.22, 12.1
for i, (lead, rest) in enumerate(items2):
    y = card_y + i * (card_h + card_gap)
    card(s2, 0.6, y, card_w, card_h, fill=CARD_BG, shadow=False, radius=0.08)
    circle_num(s2, 0.9, y + card_h/2 - 0.24, 0.48, i + 1)
    tb = s2.shapes.add_textbox(Inches(1.65), Inches(y + 0.16), Inches(card_w - 2.1), Inches(card_h - 0.3))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.line_spacing = 1.18
    r1 = p.add_run(); r1.text = lead
    r1.font.bold = True; r1.font.size = Pt(14.5); r1.font.color.rgb = MIDNIGHT; r1.font.name = FONT_BODY
    r2 = p.add_run(); r2.text = rest
    r2.font.size = Pt(14.5); r2.font.color.rgb = INK; r2.font.name = FONT_BODY

# Automated vs Still Manual legend row
leg_y = card_y + 3 * (card_h + card_gap) + 0.05
leg_w = (card_w - 0.3) / 2

pill1 = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(leg_y), Inches(leg_w), Inches(0.62))
pill1.adjustments[0] = 0.5
pill1.fill.solid(); pill1.fill.fore_color.rgb = LUCENT
no_line(pill1)
tf1 = pill1.text_frame
tf1.margin_left = Inches(0.25); tf1.vertical_anchor = MSO_ANCHOR.MIDDLE
p1 = tf1.paragraphs[0]
r = p1.add_run(); r.text = "Automated today:  "
r.font.bold = True; r.font.size = Pt(12.5); r.font.color.rgb = MIDNIGHT; r.font.name = FONT_BODY
r2 = p1.add_run(); r2.text = "analysis & root-cause classification"
r2.font.size = Pt(12.5); r2.font.color.rgb = MIDNIGHT; r2.font.name = FONT_BODY

pill2 = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6 + leg_w + 0.3), Inches(leg_y), Inches(leg_w), Inches(0.62))
pill2.adjustments[0] = 0.5
pill2.fill.solid(); pill2.fill.fore_color.rgb = AMBER_BG
no_line(pill2)
tf2 = pill2.text_frame
tf2.margin_left = Inches(0.25); tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
p2 = tf2.paragraphs[0]
r3 = p2.add_run(); r3.text = "Still manual:  "
r3.font.bold = True; r3.font.size = Pt(12.5); r3.font.color.rgb = AMBER; r3.font.name = FONT_BODY
r4 = p2.add_run(); r4.text = "exporting and uploading source files"
r4.font.size = Pt(12.5); r4.font.color.rgb = AMBER; r4.font.name = FONT_BODY


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — What's Most Important to Address: The Engineering Ask
# ══════════════════════════════════════════════════════════════════════════
s3 = add_slide()
header(s3, "What's Most Important to Address: The Engineering Ask",
       "Moving DQE from reactive and report-driven to automated.")

left_x, left_w = 0.6, 4.15
top_y = 2.32
card(s3, left_x, top_y, left_w, 3.05, fill=CARD_BG, shadow=False, radius=0.06)
txt(s3, left_x + 0.3, top_y + 0.22, left_w - 0.6, 0.3, "CURRENT PAIN POINTS", size=12, color=MIDNIGHT, bold=True, font=FONT_BODY)
pains = [
    "Little to no documentation of the DQE investigation process — difficult to standardize or automate",
    "Every investigation depends on manually exporting reports from both the PMS and Duetto",
    "No API currently exists to read (or write) the Duetto data DQE needs — this same gap blocks Auto-Provisioning and Migrations too",
    "Result: analysts manually generate and upload reports — the workflow remains largely manual",
]
bullets(s3, left_x + 0.3, top_y + 0.62, left_w - 0.6, 2.3, pains, size=11.5, color=INK,
        font=FONT_BODY, space_after=9, line_spacing=1.12, bullet_color=MIDNIGHT)

right_x, right_w = 4.95, 3.75
card(s3, right_x, top_y, right_w, 3.05, fill=CARD_BG, shadow=False, radius=0.06)
txt(s3, right_x + 0.3, top_y + 0.22, right_w - 0.6, 0.3, "WHY IT MATTERS", size=12, color=MIDNIGHT, bold=True, font=FONT_BODY)
why = [
    "Eliminates the need to manually export Duetto reports",
    "Allows DQE to retrieve required data directly",
    "Reduces manual effort and speeds up investigations",
    "Sets up the next step: a validation layer that flags anomalies automatically, rather than just replicating manual analyst steps",
]
bullets(s3, right_x + 0.3, top_y + 0.62, right_w - 0.6, 2.3, why, size=11.5, color=INK,
        font=FONT_BODY, space_after=9, line_spacing=1.12, bullet_color=MIDNIGHT)

# Far right: The Ask — Midnight card, Lucent header + tags (signature high-contrast CTA moment)
ask_x, ask_w = 8.9, 3.85
card(s3, ask_x, top_y, ask_w, 3.05, fill=MIDNIGHT, shadow=True, radius=0.06)
txt(s3, ask_x + 0.32, top_y + 0.22, ask_w - 0.64, 0.3, "THE ASK  ·  ALIGNED", size=12.5, color=LUCENT, bold=True, font=FONT_BODY)
txt(s3, ask_x + 0.32, top_y + 0.58, ask_w - 0.64, 1.15,
    "Provide API access — to read and write — for the datasets DQE needs, as part of "
    "the same API investment already agreed for Auto-Provisioning and Migrations.",
    size=12.5, color=WHITE, font=FONT_BODY, line_spacing=1.22)
txt(s3, ask_x + 0.32, top_y + 1.82, ask_w - 0.64, 0.28, "DQE DATASETS (READ)", size=10.5, color=LUCENT, bold=True, font=FONT_BODY)
ds_y = top_y + 2.14
datasets = ["Bookings", "Blocks", "Folio", "DVA"]
tag_w = (ask_w - 0.64 - 0.3) / 2
for i, ds in enumerate(datasets):
    col, row = i % 2, i // 2
    tx = ask_x + 0.32 + col * (tag_w + 0.3)
    ty = ds_y + row * 0.5
    tag = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(tx), Inches(ty), Inches(tag_w), Inches(0.4))
    tag.adjustments[0] = 0.35
    tag.fill.solid(); tag.fill.fore_color.rgb = LUCENT
    no_line(tag)
    tf = tag.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = ds
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = MIDNIGHT; r.font.name = FONT_BODY

# Bottom goal strip
goal_y = 5.62
card(s3, 0.6, goal_y, 12.1, 1.15, fill=CARD_BG, shadow=False, radius=0.08)
tb = s3.shapes.add_textbox(Inches(0.95), Inches(goal_y), Inches(11.4), Inches(1.15))
tf = tb.text_frame
tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = tf.margin_right = 0
p = tf.paragraphs[0]
p.line_spacing = 1.2
r1 = p.add_run(); r1.text = "Goal:  "
r1.font.bold = True; r1.font.size = Pt(14); r1.font.color.rgb = MIDNIGHT; r1.font.name = FONT_BODY
r2 = p.add_run()
r2.text = ("We're not asking DQE to replace analysts with AI — we're asking for the same standard, "
           "modern capability every automation project here needs. Right now we're using a postman "
           "when the ask is simply: give us email.")
r2.font.size = Pt(14); r2.font.italic = True; r2.font.color.rgb = INK; r2.font.name = FONT_BODY

prs.save("/Users/tylerlinton/DQE/DQE_Engineering_Prep_Deck.pptx")
print("saved")
